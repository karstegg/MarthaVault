#!/c/Python313/python
"""
Extract HEAL content from PowerPoint files and save as formatted text.

Supports multiple table layouts:
- Gloria: 4x2 table with Highlights/Lowlights in left/right columns
- Shafts & Winders: Multiple single-cell tables with section headers
- N3: Detected from file structure
"""

import argparse
import sys
import os
from pathlib import Path
from pptx import Presentation


def find_pptx_file(week_num, site):
    """Find the PPTX file in Week directory matching site pattern."""
    week_dir = Path(f"Week {week_num}")

    if not week_dir.exists():
        return None

    # Site-specific search patterns
    patterns = {
        "gloria": ["HEAL page*.pptx", "*gloria*.pptx"],
        "shafts-winders": ["*SHAFTS*WINDERS*.pptx", "*S&W*.pptx"],
        "n3": ["*N3*HEAL*.pptx", "*Nchwaning 3*HEAL*.pptx", "Heal N3*.pptx"],
    }

    for pattern in patterns.get(site, []):
        for pptx_file in week_dir.glob(pattern):
            return pptx_file

    return None


def extract_gloria_heal(pptx_path):
    """Extract HEAL content from Gloria's 4x2 table layout."""
    prs = Presentation(pptx_path)
    heal_data = {
        "Highlights": [],
        "Lowlights": [],
        "Emerging Issues": [],
        "Priorities": [],
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table

                # Gloria layout: 4 rows x 2 columns
                if len(table.rows) >= 4 and len(table.columns) >= 2:
                    # Row 1: Highlights (left) and Lowlights (right)
                    left_col = table.rows[1].cells[0].text.strip()
                    right_col = table.rows[1].cells[1].text.strip()

                    # Process highlights (left column) - skip if just numbering
                    if left_col and left_col not in ["1", "1.", "1,", ""]:
                        for line in left_col.split('\n'):
                            line = line.strip()
                            if line and line not in ["1", "1.", "1,"] and line not in heal_data["Highlights"]:
                                heal_data["Highlights"].append(line)

                    # Process lowlights (right column)
                    if right_col and right_col not in ["1", "1.", "1,", ""]:
                        for line in right_col.split('\n'):
                            line = line.strip()
                            if line and line not in ["1", "1.", "1,"] and line not in heal_data["Lowlights"]:
                                heal_data["Lowlights"].append(line)

                    # Row 3: Emerging Issues (left) and Priorities (right)
                    if len(table.rows) > 3:
                        left_col = table.rows[3].cells[0].text.strip()
                        right_col = table.rows[3].cells[1].text.strip()

                        # Process emerging issues - keep content even if it starts with "1."
                        if left_col and left_col not in ["", "1", "Emerging Issues"]:
                            for line in left_col.split('\n'):
                                line = line.strip()
                                # Remove leading numbering like "1." but keep the rest
                                if line.startswith("1. "):
                                    line = line[3:].strip()
                                if line and line not in heal_data["Emerging Issues"]:
                                    heal_data["Emerging Issues"].append(line)

                        # Process priorities - keep content even if it starts with "1."
                        if right_col and right_col not in ["", "1", "Priorities"]:
                            for line in right_col.split('\n'):
                                line = line.strip()
                                # Remove leading numbering like "1." but keep the rest
                                if line.startswith("1. "):
                                    line = line[3:].strip()
                                if line and line not in heal_data["Priorities"]:
                                    heal_data["Priorities"].append(line)

    return heal_data


def extract_shafts_winders_heal(pptx_path):
    """Extract HEAL content from Shafts & Winders' multi-table layout."""
    prs = Presentation(pptx_path)
    heal_data = {
        "Highlights": [],
        "Lowlights": [],
        "Emerging Issues": [],
        "Priorities": [],
        "Actions": [],
    }

    # Process only first slide (contains HEAL content, other slides are schedules)
    if len(prs.slides) > 0:
        slide = prs.slides[0]

        # Each table represents a section
        current_section = None

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                cell_text = table.rows[0].cells[0].text.strip()

                # Identify section from table content
                if cell_text.startswith("Highlights"):
                    current_section = "Highlights"
                elif cell_text.startswith("Lowlights"):
                    current_section = "Lowlights"
                elif cell_text.startswith("Emerging Issues"):
                    current_section = "Emerging Issues"
                elif cell_text.startswith("Priorities"):
                    current_section = "Priorities"
                elif cell_text.startswith("You Can Help By") or cell_text.startswith("Please provide"):
                    current_section = "Actions"

                # Extract items under this section
                if current_section:
                    for line in cell_text.split('\n'):
                        line = line.strip()
                        # Skip section headers and empty lines
                        if line and not line.endswith(":") and "Please provide" not in line:
                            if line not in heal_data[current_section]:
                                heal_data[current_section].append(line)

    return heal_data


def extract_n3_heal(pptx_path):
    """Extract HEAL content from N3 (uses Gloria-style 4x2 table layout)."""
    # N3 uses the same 4x2 table layout as Gloria
    return extract_gloria_heal(pptx_path)


def format_heal_output(heal_data):
    """Format extracted HEAL data into standard text output."""
    lines = []

    sections = ["Highlights", "Lowlights", "Emerging Issues", "Priorities", "Actions"]

    for section in sections:
        if section in heal_data:
            items = heal_data[section]

            lines.append(section)

            if items:
                for item in items:
                    # Remove numbering artifacts (e.g., "1.", "2.")
                    item = item.lstrip('0123456789. •–-')
                    item = item.strip()
                    if item:
                        lines.append(f"• {item}")

            lines.append("")  # Blank line between sections

    # Remove trailing blank lines
    while lines and lines[-1] == "":
        lines.pop()

    return "\n".join(lines)


def extract_heal(week_num, site):
    """Main extraction function."""

    # Find PPTX file
    pptx_path = find_pptx_file(week_num, site)
    if not pptx_path:
        print(f"Error: No PPTX file found for {site} in Week {week_num}", file=sys.stderr)
        return False

    print(f"Processing: {pptx_path}")

    try:
        # Extract based on site
        if site == "gloria":
            heal_data = extract_gloria_heal(pptx_path)
        elif site == "shafts-winders":
            heal_data = extract_shafts_winders_heal(pptx_path)
        elif site == "n3":
            heal_data = extract_n3_heal(pptx_path)
        else:
            print(f"Error: Unknown site: {site}", file=sys.stderr)
            return False

        # Format output
        output_text = format_heal_output(heal_data)

        # Determine output filename
        site_display = site.replace("-", " & ").title()
        output_filename = f"Week {week_num}/{site_display} HEAL Page Week{week_num}.txt"

        # Save output
        output_path = Path(output_filename)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(output_text)

        print(f"[OK] Successfully extracted HEAL content")
        print(f"[OK] Output: {output_path}")

        # Show summary
        total_items = sum(len(items) for items in heal_data.values())
        print(f"[OK] Extracted {total_items} items across {len([s for s in heal_data.values() if s])} sections")

        return True

    except Exception as e:
        print(f"Error: Failed to extract HEAL content: {e}", file=sys.stderr)
        return False


def main():
    parser = argparse.ArgumentParser(
        description="Extract HEAL content from PowerPoint files and save as formatted text"
    )

    # Support both positional and named arguments
    parser.add_argument(
        "site",
        nargs="?",
        help="Site name: gloria, shafts-winders, or n3"
    )
    parser.add_argument(
        "week",
        nargs="?",
        type=int,
        help="Week number (e.g., 16)"
    )
    parser.add_argument(
        "--site",
        dest="named_site",
        help="Site name (alternative to positional)"
    )
    parser.add_argument(
        "--week",
        type=int,
        dest="named_week",
        help="Week number (alternative to positional)"
    )

    args = parser.parse_args()

    # Determine site and week from arguments
    site = args.site or args.named_site
    week = args.week or args.named_week

    if not site or not week:
        parser.print_help()
        sys.exit(1)

    # Run extraction
    success = extract_heal(week, site)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
